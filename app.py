from flask import Flask, render_template, request, redirect, url_for,send_from_directory, flash,session , abort
import fitz  # PyMuPDF for PDFs
import pandas as pd
from docx import Document  # python-docx for DOCX
from pptx import Presentation  # python-pptx for PowerPoint
from huggingface_hub import InferenceClient
from sentence_transformers import SentenceTransformer, util
import re
import os
import pytesseract
from PIL import Image
import logging
import re
from langdetect import detect 
from transformers import AutoTokenizer, AutoModelForSequenceClassification
from torch.nn.functional import softmax
from PyPDF2 import PdfReader
from pptx import Presentation
from flask_mysqldb import MySQL
from werkzeug.utils import secure_filename
from werkzeug.security import generate_password_hash, check_password_hash
from flask import send_from_directory, current_app as app


app = Flask(__name__)
client = InferenceClient("HuggingFaceH4/zephyr-7b-beta", token="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")

#####courses-database-congif######################
# MySQL configurations
app.config['MYSQL_HOST'] = 'localhost'
app.config['MYSQL_USER'] = 'root'
app.config['MYSQL_PASSWORD'] = 'Root@nadia2024'
app.config['MYSQL_DB'] = 'pdf_database'

# Upload folder configuration
app.config['UPLOAD_FOLDER'] = r'C:/Users/pc/OneDrive\Bureau/platform/uploads'
app.config['ALLOWED_EXTENSIONS'] = {'pdf', 'docx', 'pptx', 'csv', 'png', 'jpg', 'jpeg'}

# Initialize MySQL
mysql = MySQL(app)

#secret key
app.secret_key = 'coursekey'
##################################################

# Setup logging
logging.basicConfig(level=logging.DEBUG)

# Load multilingual model for semantic analysis
model = SentenceTransformer('distiluse-base-multilingual-cased-v1')


def detect_language(text):
    try:
        return detect(text)
    except Exception as e:
        logging.error(f"Error detecting language: {str(e)}")
        return 'en'
    
 


################chatbot_related##################

# File paths for different document types (for example purposes; remove or modify as needed)
file_paths = [
    r"C:\Users\pc\Downloads\Chapitre 4, Décoissance radioactive.pdf",
    r"C:\Users\pc\Downloads\machine-readable-business.csv",

]

# Load multilingual model
model = SentenceTransformer('distiluse-base-multilingual-cased-v1')

# Initialize document embeddings
doc_embeddings = None

# Conversation history
conversation_history = []

def extract_text_from_files(file_paths):
    text = []
    for file_path in file_paths:
        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.pdf':
            doc = fitz.open(file_path)
            text.append('\n'.join([page.get_text() for page in doc]))
        elif ext == '.csv':
            df = pd.read_csv(file_path)
            text.append(df.to_string())
        elif ext == '.docx':
            doc = Document(file_path)
            text.append('\n'.join([para.text for para in doc.paragraphs]))
        elif ext == '.pptx':
            prs = Presentation(file_path)
            text.append('\n'.join([slide.shapes[0].text for slide in prs.slides if slide.shapes]))
        elif ext in ['.png', '.jpg', '.jpeg']:
            img = Image.open(file_path)
            text.append(pytesseract.image_to_string(img))
    return text

def preprocess(text):
    # Simple preprocessing example
    text = text.lower()
    text = re.sub(r'\s+', ' ', text)
    return text

def create_document_embeddings(documents):
    embeddings = model.encode(documents, convert_to_tensor=True)
    return embeddings

def retrieve_relevant_document(query, doc_embeddings):
    query_embedding = model.encode([query], convert_to_tensor=True)
    similarities = util.pytorch_cos_sim(query_embedding, doc_embeddings)[0]
    most_similar_index = similarities.argmax().item()
    return most_similar_index, similarities[most_similar_index].item()

def get_llm_response(user_message):
    response = client.chat_completion(
        [{"role": "system", "content": "You are a helpful assistant."},
         {"role": "user", "content": user_message}],
        max_tokens=1024, temperature=0.7, top_p=0.95
    )
    return response.choices[0].message.content

def generate_response(retrieved_text, user_message):
    prompt = f"Given the following text:\n{retrieved_text}\n\nAnswer the user's query:\n{user_message}"
    response = client.chat_completion(
        [{"role": "system", "content": "You are a helpful assistant."},
         {"role": "user", "content": prompt}],
        max_tokens=1024, temperature=0.7, top_p=0.95
    )
    return response.choices[0].message.content
###########################quiz-generator-relatede##########################
#########questions_generation###############################################
###1stquiz###
def generate_questions_only(text, num_questions, language=None):
    # Detect language if not provided
    if not language:
        language = detect_language(text)

    prompt = [
        {"role": "system", "content": "You are a helpful chatbot. Generate questions based on the provided context. Each question should not have any options. Provide the correct answer explicitly."},
        {"role": "user", "content": f"Context: {text}\nGenerate {num_questions} questions. Please use the same language ({language}) as the context. Indicate the correct answer for each question."}
    ]

    try:
        response = client.chat_completion(
            prompt, max_tokens=2048, temperature=0.7, top_p=0.95
        ).choices[0].message.content
        logging.debug(f"Generated response: {response}")

        # Extract and process questions and answers from the response
        questions = []
        current_question = None
        lines = response.split('\n')

        for line in lines:
            line = line.strip()
            if line.startswith('Answer:'):
                if current_question:
                    questions.append({
                        'question': current_question,
                        'answer': line[7:].strip()
                    })
                current_question = None
            elif line and not line.lower().startswith(('a)', 'b)', 'c)', 'd)')):
                if current_question:
                    questions.append({
                        'question': current_question,
                        'answer': None
                    })
                current_question = line

        if current_question:
            questions.append({
                'question': current_question,
                'answer': None
            })

        # Limit questions to the specified number
        if len(questions) > num_questions:
            questions = questions[:num_questions]

    except Exception as e:
        logging.error(f"Error generating questions only: {str(e)}")
        return []

    return questions
###2ndquiz###


def generate_questions_with_options(text, num_questions, num_options):
    language = detect_language(text)
    prompt = [
        {"role": "system", "content": "You are a helpful chatbot. Generate multiple choice questions with options and the correct answer based on the provided context. Each question should have the specified number of options. Provide the correct answer explicitly."},
        {"role": "user", "content": f"Context: {text}\nGenerate {num_questions} multiple choice questions with {num_options} options each. Please use the same language as the context. Indicate the correct answer for each question."}
    ]

    try:
        response = client.chat_completion(
            prompt, max_tokens=2048, temperature=0.7, top_p=0.95
        ).choices[0].message.content
        logging.debug(f"Generated response: {response}")

        # Extract and process questions, options, and answers from the response
        questions = []
        current_question = None
        current_options = []
        lines = response.split('\n')

        for line in lines:
            line = line.strip()
            if line.startswith('Answer:'):
                if current_question and current_options:
                    questions.append({
                        'question': current_question,
                        'options': current_options,
                        'answer': line[7:].strip()
                    })
                current_question = None
                current_options = []
            elif line and not line.lower().startswith(('a)', 'b)', 'c)', 'd)')):
                if current_question and current_options:
                    questions.append({
                        'question': current_question,
                        'options': current_options,
                        'answer': None
                    })
                current_question = line
            else:
                if current_question:
                    current_options.append(line)

        if current_question and current_options:
            questions.append({
                'question': current_question,
                'options': current_options,
                'answer': None
            })

        if len(questions) > num_questions:
            questions = questions[:num_questions]

    except Exception as e:
        logging.error(f"Error generating questions with options: {str(e)}")
        return []

    return questions


###############chatbot-routes###############################################


@app.route('/chatbot', methods=['GET', 'POST'])
def chatbot():
    global doc_embeddings
    if request.method == 'POST':
        user_message = request.form['user_input']
        file = request.files.get('file_input')
        if file:
            ext = os.path.splitext(file.filename)[1].lower()
            if ext in ['.pdf', '.csv', '.docx', '.pptx', '.png', '.jpg', '.jpeg']:
                file_path = os.path.join('uploads', file.filename)
                try:
                    file.save(file_path)
                    extracted_text = extract_text_from_files([file_path])[0]
                    preprocessed_text = preprocess(extracted_text)
                    doc_embeddings = create_document_embeddings([preprocessed_text])
                except Exception as e:
                    return render_template('app.html', bot_response=f"Error processing uploaded file: {e}", conversation_history=conversation_history)
            else:
                return render_template('app.html', bot_response="Unsupported file type.", conversation_history=conversation_history)

        if doc_embeddings is not None:
            most_similar_doc_index, similarity_score = retrieve_relevant_document(user_message, doc_embeddings)
            if similarity_score < 0.1:
                bot_response = get_llm_response(user_message)
            else:
                retrieved_text = preprocessed_text
                bot_response = generate_response(retrieved_text, user_message)
        else:
            bot_response = get_llm_response(user_message)

        conversation_history.append({'role': 'user', 'message': user_message})
        conversation_history.append({'role': 'bot', 'message': bot_response})

        return render_template('app.html', bot_response=bot_response, conversation_history=conversation_history)

    return render_template('app.html', conversation_history=conversation_history)


######login and registration#####################""
@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        email = request.form.get('email')
        password = request.form.get('password')

        # Create cursor
        cur = mysql.connection.cursor()

        try:
            # Check if user exists in the 'registration' table
            cur.execute("SELECT id, password FROM registration WHERE email = %s", (email,))
            user = cur.fetchone()
        except Exception as e:
            flash(f"Database error: {str(e)}", 'error')
            return redirect('/login')
        finally:
            cur.close()

        if user:
            user_id, hashed_password = user
            if check_password_hash(hashed_password, password):
                session['user_id'] = user_id
                return redirect('/')
            else:
                flash('Incorrect email or password', 'error')
        else:
            flash('Incorrect email or password', 'error')

    return render_template('login.html')

@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        # Extract form data
        username = request.form['username']
        name = request.form['name']
        last_name = request.form['last_name']
        email = request.form['email']
        password = request.form['password']
        hashed_password = generate_password_hash(password)
        birth_date = request.form['birth_date']

        # Create cursor
        cur = mysql.connection.cursor()

        # Check if username already exists in the 'registration' table
        cur.execute("SELECT * FROM registration WHERE username = %s", (username,))
        user_with_username = cur.fetchone()

        if user_with_username:
            flash('This username is already taken.', 'error')
            cur.close()
            return redirect('/register')

        # Check if email already exists in the 'registration' table
        cur.execute("SELECT * FROM registration WHERE email = %s", (email,))
        user_with_email = cur.fetchone()

        if user_with_email:
            flash('This email is already used.', 'error')
            cur.close()
            return redirect('/register')

        # Insert new user into the 'registration' table
        cur.execute(
            "INSERT INTO registration (username, name, last_name, email, password, birth_date) "
            "VALUES (%s, %s, %s, %s, %s, %s)",
            (username, name, last_name, email, hashed_password, birth_date)
        )
        mysql.connection.commit()
        cur.close()

        flash('Registration successful!', 'success')
        return redirect('/login')

    return render_template('register.html')

######################################
@app.route('/')
def home():
    return render_template('home.html')


###########
######################quiz-generator-routes-part####################################
###########

@app.route('/quiz')
def quiz():
    return render_template('quiz_type.html')

@app.route('/select_quiz_type', methods=['POST'])
def select_quiz_type():
    quiz_type = request.form.get('quiz_type')
    if quiz_type == 'multiple':
        return redirect(url_for('multiple_choice_quiz'))
    elif quiz_type == 'questions':
        return redirect(url_for('questions_only_quiz'))
    else:
        return redirect(url_for('quiz'))




@app.route('/multiple_choice_quiz', methods=['GET', 'POST'])
def multiple_choice_quiz():
    if request.method == 'POST':
        text = request.form['text']
        num_questions = int(request.form['num_questions'])
        num_options = int(request.form['num_options'])
        answer_type = request.form['answer_type']
        preprocessed_text = preprocess_text(text)
        quiz = generate_questions_with_options(preprocessed_text, num_questions, num_options)
        quiz_with_index = [{'index': i, **q} for i, q in enumerate(quiz)]

        # Ensure no raw data is being printed or logged directly into the response
        app.logger.debug(f"Quiz with index: {quiz_with_index}")
        
        return render_template('generated_multiple_choice_quiz.html', quiz=quiz_with_index, text=text, num_questions=num_questions, num_options=num_options, answer_type=answer_type)
    return render_template('multiple_choice_quiz.html')


@app.route('/questions_only_quiz', methods=['GET', 'POST'])
def questions_only_quiz():
    if request.method == 'POST':
        text = request.form.get('text')
        file = request.files.get('file')
        num_questions = int(request.form.get('num_questions'))

        if text:
            language = detect(text)
            # Generate quiz based on text
            questions = generate_questions_only(text, num_questions, language)
        elif file:
            file_type = file.filename.split('.')[-1].lower()
            content = ""

            if file_type == 'pdf':
                content = extract_text_from_pdf(file)
            elif file_type == 'txt':
                content = file.read().decode('utf-8')
            elif file_type == 'pptx':
                content = extract_text_from_pptx(file)

            if content:
                language = detect(content)
                # Generate quiz based on the extracted content
                questions = generate_questions_only(content, num_questions, language)
            else:
                return "Error: Could not extract content from the file."

        # Render quiz with generated questions, pass 'questions' as 'quiz'
        return render_template('questions_only_quiz_result.html', quiz=questions)

    return render_template('questions_only_quiz.html')


@app.route('/submit_questions_only_quiz', methods=['POST'])
def submit_questions_only_quiz():
    questions = request.form.getlist('question')
    correct_answers = request.form.getlist('correct_answer')
    user_answers = [request.form.get(f'answer_{i}').strip() for i in range(len(questions))]

    feedback = []

    for i in range(len(questions)):
        correct_answer = correct_answers[i]
        user_answer = user_answers[i]
        feedback.append({
            'question': questions[i],
            'correct': correct_answer,
            'user_answer': user_answer,  # Add the user's answer to the feedback
        })

    return render_template('questions_only_quiz_evaluation.html', feedback=feedback, total=len(questions))

@app.route('/submit_multiple_choice_quiz', methods=['POST'])
def submit_multiple_choice_quiz():
    answers = request.form
    feedback = []
    total_questions = int(answers.get('total_questions', 0))
    
    for i in range(total_questions):
        selected_answer = answers.getlist(f'answer_{i}')  # Get list of selected answers
        correct_answer = answers.get(f'correct_answer_{i}', None)
        question_text = answers.get(f'question_{i}', None)
        
        # Handling the case where no answer was selected
        if not selected_answer:
            selected_answer = ["None"]

        result = "Incorrect"
        if set(selected_answer) == {correct_answer}:
            result = "Correct"
        
        feedback.append({
            'index': i,
            'question': question_text,
            'selected_answer': ", ".join(selected_answer),  # Join answers if multiple
            'correct_answer': correct_answer,
            'result': result
        })
    
    score = (len([f for f in feedback if f['result'] == 'Correct']) / total_questions) * 100
    return render_template('multiple_choice_quiz_result.html', feedback=feedback, score=score)


##############courses-routes###################
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']


@app.route('/upload', methods=['GET', 'POST'])
def upload():
    if request.method == 'POST':
        coursename = request.form['coursename']
        subject = request.form['subject']
        branch = request.form['branch']
        file = request.files['file']

        if file and allowed_file(file.filename):
            # Ensure the upload directory exists
            if not os.path.exists(app.config['UPLOAD_FOLDER']):
                os.makedirs(app.config['UPLOAD_FOLDER'])

            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)

            # Insert file information into the database
            cur = mysql.connection.cursor()
            cur.execute("INSERT INTO courses (coursename, coursepath, subject, branch) VALUES (%s, %s, %s, %s)",
                        (coursename, filepath, subject, branch))
            mysql.connection.commit()
            cur.close()

            flash('File successfully uploaded')
            return redirect(url_for('home'))
    return render_template('upload.html')  # Replace with your upload page template


@app.route('/courses')
def courses():
    # Create a connection to the MySQL database
    cur = mysql.connection.cursor()
    
    # Query to get courses by subject
    subjects = [
        'Mathématiques', 'Physique', 'Chimie', 'Anglais', 
        'Sciences de la vie et de la terre', 'Philosophie'
    ]
    
    # Dictionary to hold courses by subject
    courses_by_subject = {subject: [] for subject in subjects}
    
    for subject in subjects:
        cur.execute("SELECT coursename, coursepath FROM courses WHERE subject = %s LIMIT 6", [subject])
        courses_by_subject[subject] = cur.fetchall()
    
    cur.close()
    
    return render_template('courses.html', courses_by_subject=courses_by_subject)

from flask import send_file
from werkzeug.utils import secure_filename


@app.route('/download/<filename>')
def download_file(filename):
    # Sanitize the filename
    safe_filename = secure_filename(filename)
    # Construct the file path
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], safe_filename)
    
    # Check if the file exists
    if os.path.exists(file_path):
        return send_file(file_path, as_attachment=True)

if __name__ == '__main__':
    app.run(port=5000)